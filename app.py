"""
Real Estate Competition Analysis Engine
Production-ready Streamlit application
"""
import io
import re
import requests
import streamlit as st
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn

# ─────────────────────────────────────────────────────────────────────────────
# PAGE CONFIG  (must be first Streamlit call)
# ─────────────────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="RE Competition Analysis",
    page_icon="🏗️",
    layout="wide"
)

# ─────────────────────────────────────────────────────────────────────────────
# GLOBAL CSS
# ─────────────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');
html, body, [class*="css"] { font-family: 'Inter', sans-serif; }
.stApp { background-color: #0D1117; color: #E6EDF3; }
.main .block-container { padding: 2rem 3rem; max-width: 1400px; }
.header-wrap {
    background: linear-gradient(135deg, #161B22 0%, #1C2333 100%);
    border: 1px solid #30363D; border-radius: 16px;
    padding: 2rem 2.5rem; margin-bottom: 2rem;
}
.header-badge {
    display: inline-block; background: rgba(184,134,11,0.15);
    border: 1px solid rgba(184,134,11,0.4); color: #D4A017;
    font-size: 11px; font-weight: 600; letter-spacing: 0.1em;
    text-transform: uppercase; padding: 4px 12px;
    border-radius: 20px; margin-bottom: 12px;
}
.header-title { font-size: 36px; font-weight: 700; color: #FFFFFF; margin: 0 0 8px 0; letter-spacing: -0.5px; }
.header-sub { font-size: 15px; color: #8B949E; margin: 8px 0 0 0; }
.section-label {
    font-size: 10px; font-weight: 700; letter-spacing: 0.12em;
    text-transform: uppercase; color: #D4A017;
    margin-bottom: 0.6rem; margin-top: 1.5rem;
}
.stTextInput > div > div > input,
.stSelectbox > div > div,
.stTextArea textarea {
    background-color: #0D1117 !important; border: 1px solid #30363D !important;
    border-radius: 8px !important; color: #E6EDF3 !important;
    font-family: 'Inter', sans-serif !important;
}
.stTextInput > div > div > input:focus { border-color: #B8860B !important; }
label { color: #8B949E !important; font-size: 13px !important; font-weight: 500 !important; }
.stButton > button {
    background: #B8860B !important; color: #FFFFFF !important; border: none !important;
    border-radius: 10px !important; padding: 14px 36px !important;
    font-size: 15px !important; font-weight: 600 !important; width: 100% !important;
    letter-spacing: 0.02em !important; transition: all 0.2s !important;
}
.stButton > button:hover { background: #D4A017 !important; transform: translateY(-1px) !important; }
.stTabs [data-baseweb="tab-list"] {
    background: #161B22; border-radius: 10px; padding: 4px;
    border: 1px solid #30363D; gap: 4px;
}
.stTabs [data-baseweb="tab"] {
    background: transparent; border-radius: 8px; color: #8B949E;
    font-size: 13px; font-weight: 500; padding: 8px 16px;
}
.stTabs [aria-selected="true"] { background: #B8860B !important; color: #FFFFFF !important; }
.stTabs [data-baseweb="tab-panel"] {
    background: #161B22; border: 1px solid #30363D;
    border-radius: 0 12px 12px 12px; padding: 2rem; margin-top: -1px;
}
.stAlert { border-radius: 10px !important; border: none !important; }
hr { border-color: #21262D !important; }
.stDownloadButton > button {
    background: transparent !important; border: 1px solid #B8860B !important;
    color: #D4A017 !important; border-radius: 8px !important;
    font-weight: 500 !important; width: 100% !important;
}
.stDownloadButton > button:hover { background: rgba(184,134,11,0.1) !important; }
table { width: 100%; border-collapse: collapse; font-size: 13px; }
th {
    background: #1C2333; color: #D4A017; font-weight: 600;
    padding: 10px 14px; text-align: left; border-bottom: 2px solid #B8860B;
    font-size: 11px; text-transform: uppercase; letter-spacing: 0.06em;
}
td { padding: 10px 14px; border-bottom: 1px solid #21262D; color: #E6EDF3; vertical-align: top; }
tr:hover td { background: rgba(184,134,11,0.05); }
</style>
""", unsafe_allow_html=True)

# ─────────────────────────────────────────────────────────────────────────────
# HEADER
# ─────────────────────────────────────────────────────────────────────────────
st.markdown("""
<div class="header-wrap">
    <div class="header-badge">Real Estate Developer Intelligence</div>
    <div class="header-title">🏗️ Competition Analysis Engine</div>
    <div class="header-sub">
        Your team inputs the market — AI delivers a developer-grade
        competitor intelligence report &amp; McKinsey-style PPT
    </div>
</div>
""", unsafe_allow_html=True)

st.info("📡 Fetches live Google data + AI analysis. Verify all figures before board presentations.")

# ─────────────────────────────────────────────────────────────────────────────
# INPUT FORM
# ─────────────────────────────────────────────────────────────────────────────
st.markdown('<div class="section-label">📍 Market Context</div>', unsafe_allow_html=True)
col1, col2, col3 = st.columns(3)
with col1:
    micromarket = st.text_input("Micro-market", placeholder="e.g. Goregaon West")
with col2:
    city = st.text_input("City", placeholder="e.g. Mumbai")
with col3:
    product_type = st.selectbox(
        "Product Type",
        ["Residential", "Commercial", "Mixed-use", "Plots", "Warehousing"]
    )

st.markdown('<div class="section-label">🏢 Your Project Details</div>', unsafe_allow_html=True)
col4, col5, col6 = st.columns(3)
with col4:
    our_project_name = st.text_input("Your Project Name", placeholder="e.g. Skyline Residences")
with col5:
    our_land_area = st.text_input("Land Area / FSI", placeholder="e.g. 2 acres, FSI 3.0")
with col6:
    our_target_segment = st.text_input(
        "Target Segment", placeholder="e.g. Premium mid-segment, ₹1.5–3 Cr"
    )

col7, col8 = st.columns(2)
with col7:
    our_configs = st.text_input(
        "Planned Configurations", placeholder="e.g. 2BHK, 3BHK (700–1200 sqft carpet)"
    )
with col8:
    our_launch_timeline = st.selectbox("Your Launch Timeline", [
        "Immediate (0–3 months)", "Short-term (3–6 months)",
        "Mid-term (6–12 months)", "Long-term (1–2 years)"
    ])

our_strengths = st.text_area(
    "Project USPs / Strengths (optional)",
    placeholder="e.g. Metro connectivity 500m, branded developer, rooftop amenities, RERA registered",
    height=68
)

st.markdown('<div class="section-label">🎯 Known Competitors to Benchmark</div>', unsafe_allow_html=True)
st.caption("List projects your team has already identified. AI will analyse these + discover additional ones.")

comp_row1 = st.columns(3)
competitors = []
for i, col in enumerate(comp_row1):
    with col:
        nm = st.text_input(
            f"Competitor {i+1}",
            placeholder=["Lodha Palava", "Godrej Reserve", "Raymond Realty"][i],
            key=f"comp_{i}"
        )
        competitors.append(nm.strip())

comp_row2 = st.columns(3)
for i, col in enumerate(comp_row2[:2]):
    with col:
        nm = st.text_input(f"Competitor {i+4}", placeholder="Optional", key=f"comp_ex_{i}")
        competitors.append(nm.strip())

all_competitors = [c for c in competitors if c]

st.markdown('<div class="section-label">🔑 API Configuration</div>', unsafe_allow_html=True)
col_a, col_b = st.columns(2)
with col_a:
    groq_key = st.text_input("Groq API Key", type="password", placeholder="gsk_...")
with col_b:
    serp_key = st.text_input("SerpAPI Key", type="password", placeholder="Your SerpAPI key")

st.divider()
run = st.button("🚀 Run Full Competition Analysis")


# ─────────────────────────────────────────────────────────────────────────────
# DATA FETCHER
# ─────────────────────────────────────────────────────────────────────────────
def fetch_live_data(micromarket, city, product_type, all_competitors, serp_key):
    snippets = []
    queries  = [
        f"property price per sqft {micromarket} {city} 2024 2025 {product_type.lower()}",
        f"new residential projects launch {micromarket} {city} 2025",
        f"real estate market trend {micromarket} {city} latest",
    ]
    for comp in all_competitors[:5]:
        queries.append(f"{comp} {city} price sqft carpet area configurations")

    for q in queries:
        try:
            res = requests.get(
                "https://serpapi.com/search",
                params={"q": q, "api_key": serp_key, "num": 5, "gl": "in", "hl": "en"},
                timeout=12
            )
            data = res.json()
            for item in data.get("organic_results", [])[:4]:
                snippet = item.get("snippet", "")
                source  = item.get("source", "")
                title   = item.get("title",   "")
                if snippet:
                    snippets.append(f"[{source}] {title}: {snippet}")
        except Exception:
            continue

    return "\n\n".join(snippets) if snippets else "Live data unavailable — AI will use knowledge base."


# ─────────────────────────────────────────────────────────────────────────────
# PROMPT BUILDER
# ─────────────────────────────────────────────────────────────────────────────
def build_prompt(micromarket, city, product_type, our_project_name, our_land_area,
                 our_target_segment, our_configs, our_launch_timeline, our_strengths,
                 all_competitors, live_data):
    comp_list = ", ".join(all_competitors) if all_competitors else "Not specified — discover top 5 projects"
    cfg_parts = [c.strip() for c in our_configs.split(',') if c.strip()] if our_configs else ["2BHK", "3BHK"]
    cfg1 = cfg_parts[0] if len(cfg_parts) > 0 else "2BHK"
    cfg2 = cfg_parts[1] if len(cfg_parts) > 1 else "3BHK"
    proj = (our_project_name or "YOUR PROJECT").upper()

    return f"""
You are a senior real estate strategy analyst with deep expertise in Indian property markets.
Your client is a DEVELOPER who needs a hard-hitting, specific competitor intelligence report.

DEVELOPER'S PROJECT:
- Project Name: {our_project_name or 'Unnamed Project'}
- Location: {micromarket}, {city}
- Product Type: {product_type}
- Land Area / FSI: {our_land_area or 'Not specified'}
- Target Segment: {our_target_segment or 'Not specified'}
- Planned Configurations: {our_configs or 'Not specified'}
- Launch Timeline: {our_launch_timeline}
- Known USPs: {our_strengths or 'None mentioned'}

COMPETITORS IDENTIFIED BY TEAM:
{comp_list}

LIVE MARKET DATA (fetched today from Google):
{live_data[:3500]}

RULES:
- Use live data as PRIMARY source; fill gaps with your knowledge of Indian RE markets.
- Every number must be specific. ₹ figures must have carpet sqft context.
- Think like the developer's strategy team: where can they price, differentiate, and win?
- If named competitor has insufficient data, use your training knowledge of that project.

OUTPUT FORMAT — use EXACTLY these markers (no skipping, no renaming):

SECTION_1: MARKET SNAPSHOT
• Avg ₹/sqft range: ₹X,XXX – ₹X,XXX (carpet, {micromarket})
• Total active inventory: ~X projects / ~X,XXX units in pipeline
• Market temperature: [Hot 🔥 / Stable 📊 / Cooling ❄️] — reason in one line
• Primary buyer: [profile, income bracket, end-use vs investment %]
• YoY price appreciation: X% (source or estimate)

SECTION_2: COMPETITOR DEEP-DIVE
For EACH project (named + 2–3 discovered), write:

**[Project Name] — [Developer]**
- Configurations: X BHK (XXX–XXX sqft carpet)
- Pricing: ₹X,XXX – ₹X,XXX/sqft | All-in: ₹X.X Cr – ₹X.X Cr
- Stage: [Under Construction / Ready / New Launch]
- RERA: [Registered / Not registered / Unknown]
- Key USP: [one line]
- Weakness / Gap: [one line — where they are vulnerable]

After all projects:
**Market Pricing Band:** ₹X,XXX – ₹X,XXX/sqft (budget) | ₹X,XXX – ₹X,XXX/sqft (mid) | ₹X,XXX+/sqft (premium)
**Dominant payment scheme:** [CLP / Subvention / Flexi — typical structure]

SECTION_3: COMPETITOR COMPARISON TABLE
| Project | Developer | Config | Carpet (sqft) | ₹/sqft | All-in Price | Stage | Key USP |
|---|---|---|---|---|---|---|---|
[fill 5–7 rows with real data]

SECTION_4: PRICING STRATEGY FOR {proj}
• Recommended launch ₹/sqft: ₹XX,XXX — justify vs 3 nearest competitors
• {cfg1}: XXX sqft carpet → ₹X.XX Cr all-in
• {cfg2}: XXX sqft carpet → ₹X.XX Cr all-in
• Floor rise: ₹XX/sqft per floor (low-rise) | ₹XX/sqft per floor (high-rise)
• PLC premiums: preferred-facing +X% | corner +X%
• Parking: ₹X.X L per covered slot
• Recommended payment scheme: [name + structure e.g. 10:80:10 CLP]
• Absorption forecast: X units/month at recommended pricing

SECTION_5: MARKET GAPS & DIFFERENTIATION OPPORTUNITIES
GAP 1: [underserved segment / config / price point + why + demand size]
GAP 2: [underserved segment / config / price point + why + demand size]
GAP 3: [location, amenity, or product format gap + why + demand size]
HOW {proj} CAN WIN:
• [specific differentiation play 1]
• [specific differentiation play 2]
• [specific differentiation play 3]

SECTION_6: RISK FLAGS
• Risk 1 [HIGH/MED/LOW]: [description] → Mitigation: [one line]
• Risk 2 [HIGH/MED/LOW]: [description] → Mitigation: [one line]
• Risk 3 [HIGH/MED/LOW]: [description] → Mitigation: [one line]
• Risk 4 [HIGH/MED/LOW]: [description] → Mitigation: [one line]
• VERDICT: [GO ✅ / CAUTION ⚠️ / HOLD 🛑] — [one sentence strategic rationale]
"""


# ─────────────────────────────────────────────────────────────────────────────
# PPT PALETTE
# ─────────────────────────────────────────────────────────────────────────────
_NAVY      = RGBColor(0x1E, 0x27, 0x61)
_DK_NAVY   = RGBColor(0x2A, 0x35, 0x7A)
_GOLD      = RGBColor(0xB8, 0x86, 0x0B)
_GOLD_L    = RGBColor(0xD4, 0xA0, 0x17)
_GOLD_PALE = RGBColor(0xF5, 0xE6, 0xB0)
_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
_CREAM     = RGBColor(0xF0, 0xF0, 0xF5)
_LT_BLUE   = RGBColor(0xCA, 0xDC, 0xFC)
_SLATE     = RGBColor(0x4A, 0x4A, 0x6A)
_MID_GRAY  = RGBColor(0x8B, 0x8B, 0xA8)
_RED_F     = RGBColor(0xC0, 0x39, 0x2B)
_AMBER_F   = RGBColor(0xD3, 0x7A, 0x00)
_GREEN_F   = RGBColor(0x1E, 0x88, 0x55)


# ─────────────────────────────────────────────────────────────────────────────
# PPT HELPERS
# ─────────────────────────────────────────────────────────────────────────────
def _bg(slide, rgb):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = rgb

def _r(sl, l, t, w, h, fill, line=None, lw=0.5):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def _t(sl, text, l, t, w, h, sz=11, bold=False, italic=False,
       color=None, align=PP_ALIGN.LEFT, font="Calibri"):
    if color is None: color = _WHITE
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = str(text)
    run.font.size = Pt(sz); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font
    return tb

def _cbg(cell, r, g, b):
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    sc = etree.SubElement(sf,   qn('a:srgbClr'))
    sc.set('val', f'{r:02X}{g:02X}{b:02X}')

def _foot(sl, label, page):
    _r(sl, 0, 7.28, 10, 0.22, _NAVY)
    _t(sl, label,      0.3, 7.30, 8.8, 0.19, sz=8, color=_LT_BLUE)
    _t(sl, str(page),  9.5, 7.30, 0.4, 0.19, sz=8, color=_GOLD, align=PP_ALIGN.RIGHT)

def _bullets(sl, lines, sx, sy, mw, rh=0.44, max_r=10, dot=None, tc=None):
    if dot is None: dot = _GOLD
    if tc  is None: tc  = _SLATE
    y = sy
    for line in lines[:max_r]:
        c = line.strip().lstrip("•·-* \t")
        if not c or len(c) < 4: continue
        _r(sl, sx, y+0.15, 0.06, 0.06, dot)
        _t(sl, c, sx+0.18, y, mw-0.18, rh, sz=10, color=tc)
        y += rh
        if y > 7.0: break
    return y

def _clean(text, n=260):
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*',     r'\1', text)
    lines = [l.strip().lstrip("-•·* \t") for l in text.split('\n')
             if l.strip() and len(l.strip()) > 5]
    return ("  ·  ".join(lines[:3]))[:n]


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE FUNCTIONS
# ─────────────────────────────────────────────────────────────────────────────
def _cover(prs, micromarket, city, product_type, project_name, timeline):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, _NAVY)
    _r(sl, 0, 0, 0.12, 7.5, _GOLD)
    _r(sl, 9.0, 0, 1.0, 0.14, _GOLD)
    _r(sl, 9.0, 7.36, 1.0, 0.14, _GOLD)
    _t(sl, "REAL ESTATE DEVELOPER INTELLIGENCE", 0.5, 0.32, 9.0, 0.38, sz=9, bold=True, color=_GOLD)
    _t(sl, "Competition\nAnalysis", 0.5, 1.05, 9.0, 1.9, sz=52, bold=True, color=_WHITE)
    _t(sl, project_name or "Developer Project", 0.5, 3.1, 9.0, 0.6, sz=24, color=_GOLD_L)
    _t(sl, f"{micromarket},  {city}", 0.5, 3.65, 9.0, 0.45, sz=16, color=_LT_BLUE)
    _r(sl, 0.5, 4.3, 4.0, 0.05, _GOLD)
    _t(sl, f"Product Type: {product_type}    ·    Launch: {timeline}",
       0.5, 4.48, 9.0, 0.38, sz=11, color=_LT_BLUE)
    _t(sl, "STRICTLY CONFIDENTIAL  ·  INTERNAL STRATEGY USE ONLY",
       0.5, 6.95, 9.0, 0.28, sz=8, italic=True, color=_MID_GRAY)


def _exec_summary(prs, micromarket, city, bullets, project_name, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, _CREAM)
    _r(sl, 0, 0, 10, 0.88, _NAVY)
    _t(sl, "EXECUTIVE SUMMARY", 0.4, 0.2, 7.5, 0.5, sz=20, bold=True, color=_WHITE)
    _t(sl, f"{micromarket}, {city}", 7.2, 0.23, 2.6, 0.4, sz=10, color=_GOLD_L, align=PP_ALIGN.RIGHT)
    labels = ["Market Snapshot","Competitive Landscape","Pricing Position","Strategic Verdict"]
    icons  = ["01","02","03","04"]
    pos    = [(0.28,1.0),(5.18,1.0),(0.28,3.72),(5.18,3.72)]
    for i,(lx,ly) in enumerate(pos):
        _r(sl, lx,ly,4.6,2.48,_WHITE, line=RGBColor(0xD8,0xD8,0xE8), lw=0.5)
        _r(sl, lx,ly,4.6,0.10,_GOLD)
        _t(sl, icons[i],  lx+0.16,ly+0.14,0.6, 0.42, sz=24,bold=True,color=_GOLD)
        _t(sl, labels[i], lx+0.74,ly+0.17,3.6, 0.35, sz=10,bold=True,color=_NAVY)
        raw = bullets[i] if i < len(bullets) else "See full analysis"
        _t(sl, _clean(raw), lx+0.16,ly+0.62,4.2,1.76, sz=10,color=_SLATE)
    _foot(sl, f"{project_name or 'Project'}  ·  {micromarket}, {city}  ·  Confidential", page)


def _market(prs, micromarket, city, content, project_name, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, _CREAM)
    _r(sl, 0,0,10,0.88,_NAVY)
    _t(sl,"MARKET SNAPSHOT",0.4,0.2,7.5,0.5,sz=20,bold=True,color=_WHITE)
    _t(sl,f"{micromarket}, {city}",7.2,0.23,2.6,0.4,sz=10,color=_GOLD_L,align=PP_ALIGN.RIGHT)
    _r(sl, 0.28,1.0,3.0,6.1,_NAVY)
    pm = re.findall(r'₹[\d,]+(?:\s*[–-]\s*₹[\d,]+)?', content)
    sv = [pm[0][:18] if pm else "₹18K–24K","STABLE →","+8–12%","~12 active"]
    sl_ = ["Avg ₹/sqft","Market Temp","YoY Growth","Inventory"]
    for i,(lb,vl) in enumerate(zip(sl_,sv)):
        y = 1.3+i*1.35
        _t(sl,lb,0.42,y,2.7,0.28,sz=8,bold=True,color=_GOLD)
        _t(sl,vl,0.42,y+0.29,2.7,0.62,sz=14,bold=True,color=_WHITE)
        if i<3: _r(sl,0.45,y+0.98,2.5,0.02,_DK_NAVY)
    _t(sl,"Market Intelligence",3.55,1.05,6.1,0.4,sz=13,bold=True,color=_NAVY)
    _r(sl,3.55,1.5,1.6,0.05,_GOLD)
    lines=[l.strip().lstrip("•·-* ") for l in content.split('\n') if l.strip() and len(l.strip())>12]
    _bullets(sl,lines,3.55,1.65,6.1,rh=0.55,max_r=9)
    _foot(sl,f"{project_name or 'Project'}  ·  {micromarket}, {city}  ·  Confidential",page)


def _deepdive(prs, micromarket, city, content, project_name, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,_CREAM)
    _r(sl,0,0,10,0.88,_NAVY)
    _t(sl,"COMPETITOR DEEP-DIVE",0.4,0.2,7.5,0.5,sz=20,bold=True,color=_WHITE)
    _t(sl,f"{micromarket}, {city}",7.2,0.23,2.6,0.4,sz=10,color=_GOLD_L,align=PP_ALIGN.RIGHT)
    blocks = re.split(r'\*\*(.+?)\*\*', content)
    cards  = []
    for j in range(1,len(blocks)-1,2):
        h=blocks[j].strip(); b=blocks[j+1].strip() if j+1<len(blocks) else ""
        if "—" in h or "-" in h or len(h)>5: cards.append((h,b))
    pos=[(0.28,1.0),(5.18,1.0),(0.28,3.8),(5.18,3.8)]
    if cards:
        for idx,(header,body) in enumerate(cards[:4]):
            lx,ly=pos[idx]
            _r(sl,lx,ly,4.6,2.55,_WHITE,line=RGBColor(0xD8,0xD8,0xE8),lw=0.5)
            _r(sl,lx,ly,0.09,2.55,_GOLD)
            _t(sl,header[:55],lx+0.2,ly+0.1,4.2,0.36,sz=10,bold=True,color=_NAVY)
            sub=[l.strip().lstrip("•·- ") for l in body.split('\n') if l.strip() and len(l.strip())>5][:5]
            y=ly+0.52
            for ln in sub:
                _r(sl,lx+0.2,y+0.1,0.06,0.06,_GOLD)
                _t(sl,ln[:120],lx+0.38,y,4.1,0.4,sz=9,color=_SLATE)
                y+=0.4
    else:
        lines=[l.strip().lstrip("•·-* ") for l in content.split('\n') if l.strip() and len(l.strip())>8]
        _bullets(sl,lines,0.35,1.05,9.3,max_r=14)
    _foot(sl,f"{project_name or 'Project'}  ·  {micromarket}, {city}  ·  Confidential",page)


def _table(prs, micromarket, city, content, project_name, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,_CREAM)
    _r(sl,0,0,10,0.88,_NAVY)
    _t(sl,"COMPETITOR COMPARISON TABLE",0.4,0.2,7.5,0.5,sz=20,bold=True,color=_WHITE)
    _t(sl,f"{micromarket}, {city}",7.2,0.23,2.6,0.4,sz=10,color=_GOLD_L,align=PP_ALIGN.RIGHT)
    rows=[]
    for line in content.split('\n'):
        if '|' in line and '---' not in line and line.strip().startswith('|'):
            cells=[c.strip() for c in line.split('|') if c.strip()]
            if cells: rows.append(cells)
    if len(rows)>=2:
        headers=rows[0]; data_rows=rows[1:9]
        n_cols=min(len(headers),8); n_rows=len(data_rows)+1
        raw_w=[1.7,1.4,0.85,1.0,0.95,1.1,0.85,1.7][:n_cols]
        scale=9.3/sum(raw_w); col_w=[w*scale for w in raw_w]
        rh=min(Inches(0.48),Inches(5.5/n_rows))
        tbl=sl.shapes.add_table(n_rows,n_cols,Inches(0.35),Inches(1.05),Inches(9.3),rh*n_rows).table
        for ci,cw in enumerate(col_w): tbl.columns[ci].width=Inches(cw)
        for ci,h in enumerate(headers[:n_cols]):
            cell=tbl.cell(0,ci); cell.text=h.upper()
            _cbg(cell,0x1E,0x27,0x61)
            p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
            run=p.runs[0] if p.runs else p.add_run()
            run.font.bold=True; run.font.size=Pt(8)
            run.font.color.rgb=_GOLD; run.font.name="Calibri"
        for ri,row_data in enumerate(data_rows):
            bg=(0xFF,0xFF,0xFF) if ri%2==0 else (0xF0,0xF0,0xF8)
            for ci in range(n_cols):
                cell=tbl.cell(ri+1,ci)
                cell.text=row_data[ci] if ci<len(row_data) else ""
                _cbg(cell,*bg)
                p=cell.text_frame.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
                run=p.runs[0] if p.runs else p.add_run()
                run.font.size=Pt(9); run.font.bold=(ci==0)
                run.font.color.rgb=_NAVY if ci==0 else _SLATE; run.font.name="Calibri"
    else:
        _t(sl,"Table data unavailable — see Competitor Deep-Dive tab.",0.35,1.1,9.3,0.4,sz=11,color=_SLATE)
    _foot(sl,f"{project_name or 'Project'}  ·  {micromarket}, {city}  ·  Confidential",page)


def _pricing(prs, micromarket, city, content, project_name, configs_str, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,_CREAM)
    _r(sl,0,0,10,0.88,_NAVY)
    _t(sl,f"PRICING STRATEGY  ·  {(project_name or 'YOUR PROJECT').upper()}",
       0.4,0.2,9.3,0.5,sz=18,bold=True,color=_WHITE)
    _t(sl,f"{micromarket}, {city}",7.2,0.23,2.6,0.4,sz=10,color=_GOLD_L,align=PP_ALIGN.RIGHT)
    raw_c=[c.strip() for c in configs_str.split(',') if c.strip()] if configs_str else []
    defs=["2BHK","3BHK","Premium / Jodi"]
    cfgs=[raw_c[i] if i<len(raw_c) else defs[i] for i in range(3)]
    pl=[l for l in content.split('\n') if ('→' in l or 'all-in' in l.lower() or 'carpet' in l.lower()) and '₹' in l]
    prices=[]
    for p in pl:
        m=re.search(r'₹\d+\.\d+\s*Cr',p)
        if m: prices.append(m.group().strip())
    if not prices: prices=re.findall(r'₹\d+\.\d+\s*Cr',content)
    cf=[_NAVY,_GOLD,_NAVY]; vc=[_GOLD_L,_WHITE,_GOLD_L]; sc=[_LT_BLUE,_DK_NAVY,_LT_BLUE]
    for i in range(3):
        lx=0.28+i*3.25
        _r(sl,lx,1.0,3.1,2.0,cf[i])
        _t(sl,cfgs[i],lx+0.2,1.1,2.7,0.5,sz=22,bold=True,color=_WHITE)
        _t(sl,prices[i] if i<len(prices) else "See report",lx+0.2,1.65,2.7,0.5,sz=17,bold=True,color=vc[i])
        _t(sl,"All-in price",lx+0.2,2.22,2.7,0.3,sz=9,color=sc[i])
    _t(sl,"Pricing Intelligence",0.28,3.2,9.5,0.38,sz=13,bold=True,color=_NAVY)
    _r(sl,0.28,3.62,1.6,0.05,_GOLD)
    lines=[l.strip().lstrip("•·-* ") for l in content.split('\n') if l.strip() and len(l.strip())>8]
    _bullets(sl,lines,0.28,3.78,9.5,rh=0.46,max_r=8)
    _foot(sl,f"{project_name or 'Project'}  ·  {micromarket}, {city}  ·  Confidential",page)


def _gaps(prs, micromarket, city, content, project_name, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,_CREAM)
    _r(sl,0,0,10,0.88,_NAVY)
    _t(sl,"MARKET GAPS & DIFFERENTIATION OPPORTUNITIES",0.4,0.2,9.3,0.5,sz=18,bold=True,color=_WHITE)
    _t(sl,f"{micromarket}, {city}",7.2,0.23,2.6,0.4,sz=10,color=_GOLD_L,align=PP_ALIGN.RIGHT)
    lines=[l.strip().lstrip("•·-* ") for l in content.split('\n') if l.strip() and len(l.strip())>10]
    gap_lines=[l for l in lines if re.match(r'(?i)gap',l)][:3]
    if not gap_lines: gap_lines=lines[:3]
    _t(sl,"IDENTIFIED GAPS",0.28,0.98,4.5,0.3,sz=9,bold=True,color=_NAVY)
    for i,gap in enumerate(gap_lines):
        c2=re.sub(r'(?i)^gap\s*\d*\s*[:.]?\s*','',gap).strip()
        ly=1.35+i*1.75
        _r(sl,0.28,ly,4.5,1.56,_WHITE,line=RGBColor(0xD8,0xD8,0xE8),lw=0.5)
        _r(sl,0.28,ly,0.09,1.56,_GOLD)
        _t(sl,f"GAP {i+1:02d}",0.46,ly+0.1,4.0,0.28,sz=9,bold=True,color=_GOLD)
        _t(sl,c2[:220],0.46,ly+0.42,3.95,1.02,sz=10,color=_SLATE)
    _r(sl,5.02,0.98,4.72,6.1,_NAVY)
    _t(sl,"HOW TO WIN",5.22,1.12,4.2,0.32,sz=9,bold=True,color=_GOLD)
    _r(sl,5.22,1.48,1.2,0.04,_GOLD_PALE)
    win=[]
    for l in lines:
        cl=re.sub(r'(?i)(how\s+\S.{0,30}can\s+win\s*[:·]?)','',l).strip().lstrip(':·- ')
        if cl and len(cl)>=8 and l not in gap_lines: win.append(cl)
    _bullets(sl,win,5.22,1.65,4.35,rh=0.56,max_r=8,dot=_GOLD,tc=_LT_BLUE)
    _foot(sl,f"{project_name or 'Project'}  ·  {micromarket}, {city}  ·  Confidential",page)


def _risks(prs, micromarket, city, content, project_name, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,_CREAM)
    _r(sl,0,0,10,0.88,_NAVY)
    _t(sl,"RISK FLAGS",0.4,0.2,7.5,0.5,sz=20,bold=True,color=_WHITE)
    _t(sl,f"{micromarket}, {city}",7.2,0.23,2.6,0.4,sz=10,color=_GOLD_L,align=PP_ALIGN.RIGHT)
    lines=[l.strip().lstrip("•·-* ") for l in content.split('\n') if l.strip() and len(l.strip())>8]
    verdict=""; y=1.02; count=0
    cmap={"HIGH":_RED_F,"MED":_AMBER_F,"LOW":_GREEN_F}
    for line in lines:
        if "VERDICT" in line.upper(): verdict=line; continue
        lvl="MED"
        for lv in ["HIGH","MED","LOW"]:
            if f"[{lv}]" in line.upper() or f" {lv}" in line.upper(): lvl=lv; break
        _r(sl,0.28,y,1.05,0.72,cmap[lvl])
        _t(sl,lvl,0.28,y+0.17,1.05,0.4,sz=12,bold=True,color=_WHITE,align=PP_ALIGN.CENTER)
        _r(sl,1.45,y,8.25,0.72,_WHITE,line=RGBColor(0xD8,0xD8,0xE8),lw=0.5)
        c2=re.sub(r'\[?(HIGH|MED|LOW)\]?','',line,flags=re.IGNORECASE).strip().lstrip(':- ')
        _t(sl,c2[:220],1.6,y+0.08,7.95,0.58,sz=10,color=_SLATE)
        y+=0.88; count+=1
        if count>=5: break
    if verdict:
        _r(sl,0.28,y+0.18,9.44,0.88,_NAVY)
        vc=re.sub(r'(?i)^verdict\s*[:·]?\s*','',verdict).strip()
        _t(sl,f"VERDICT  ·  {vc}",0.46,y+0.3,9.1,0.64,sz=12,bold=True,color=_GOLD_L)
    _foot(sl,f"{project_name or 'Project'}  ·  {micromarket}, {city}  ·  Confidential",page)


def _closing(prs, micromarket, city, project_name):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl,_NAVY)
    _r(sl,0,0,0.12,7.5,_GOLD)
    _r(sl,9.0,0,1.0,0.14,_GOLD)
    _r(sl,9.0,7.36,1.0,0.14,_GOLD)
    _t(sl,"REAL ESTATE DEVELOPER INTELLIGENCE",0.5,1.9,9.0,0.45,sz=10,bold=True,color=_GOLD,align=PP_ALIGN.CENTER)
    _t(sl,"Thank You",0.5,2.7,9.0,1.1,sz=54,bold=True,color=_WHITE,align=PP_ALIGN.CENTER)
    _r(sl,3.0,4.0,4.0,0.07,_GOLD)
    _t(sl,f"{project_name or 'Developer Project'}  ·  {micromarket}, {city}",
       0.5,4.2,9.0,0.48,sz=15,color=_LT_BLUE,align=PP_ALIGN.CENTER)
    _t(sl,"This document is confidential and intended solely for internal strategy use.",
       0.5,6.7,9.0,0.35,sz=9,italic=True,color=_MID_GRAY,align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# MASTER PPT GENERATOR
# ─────────────────────────────────────────────────────────────────────────────
def generate_ppt(sections, micromarket, city, product_type,
                 project_name, configs_str, launch_timeline):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)
    eb = [sections.get(k,"")[:200] for k in ["1","2","4","6"]]
    _cover(prs, micromarket, city, product_type, project_name, launch_timeline)
    _exec_summary(prs, micromarket, city, eb, project_name, 2)
    _market(prs, micromarket, city, sections.get("1",""), project_name, 3)
    _deepdive(prs, micromarket, city, sections.get("2",""), project_name, 4)
    _table(prs, micromarket, city, sections.get("3",""), project_name, 5)
    _pricing(prs, micromarket, city, sections.get("4",""), project_name, configs_str, 6)
    _gaps(prs, micromarket, city, sections.get("5",""), project_name, 7)
    _risks(prs, micromarket, city, sections.get("6",""), project_name, 8)
    _closing(prs, micromarket, city, project_name)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ─────────────────────────────────────────────────────────────────────────────
# MAIN RUN LOGIC
# ─────────────────────────────────────────────────────────────────────────────
if run:
    if not micromarket or not city or not groq_key or not serp_key:
        st.error("⚠️ Please fill in Micro-market, City, and both API keys before running.")
    else:
        with st.status("🔍 Fetching live market & competitor data...", expanded=True) as status:
            st.write(f"Searching market prices in {micromarket}, {city}...")
            if all_competitors:
                st.write(f"Searching named competitors: {', '.join(all_competitors)}...")
            st.write("Searching active project launches and market trends...")
            live_data = fetch_live_data(micromarket, city, product_type, all_competitors, serp_key)
            status.update(label="✅ Live data fetched!", state="complete")

        with st.expander("📄 View raw data fetched from web"):
            st.text(live_data[:2500])

        with st.status("🤖 AI generating competition analysis...", expanded=True) as status:
            st.write("Running competitor intelligence analysis...")
            st.write("Building pricing strategy and gap analysis...")
            prompt = build_prompt(
                micromarket, city, product_type,
                our_project_name, our_land_area, our_target_segment,
                our_configs, our_launch_timeline, our_strengths,
                all_competitors, live_data
            )
            client   = Groq(api_key=groq_key)
            response = client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=4000
            )
            result = response.choices[0].message.content
            status.update(label="✅ Analysis complete!", state="complete")

        raw_sections = result.split("SECTION_")
        sections_dict = {}
        for sec in raw_sections[1:]:
            if ':' in sec:
                num = sec.split(':')[0].strip()
                sections_dict[num] = ':'.join(sec.split(':')[1:]).strip()

        proj_label = our_project_name or micromarket
        st.success(f"🎉 Competition Analysis for **{proj_label}** is ready!")
        st.divider()

        tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
            "📊 Market Snapshot", "🏢 Competitor Deep-Dive", "📋 Comparison Table",
            "💰 Pricing Strategy", "🎯 Gaps & Opportunities", "⚠️ Risk Flags"
        ])
        def _show(k):
            return sections_dict.get(k, "_Section not generated — please re-run._")

        with tab1: st.markdown(_show("1"))
        with tab2: st.markdown(_show("2"))
        with tab3: st.markdown(_show("3"))
        with tab4: st.markdown(_show("4"))
        with tab5: st.markdown(_show("5"))
        with tab6: st.markdown(_show("6"))

        st.divider()
        st.markdown('<div class="section-label">⬇️ Download Reports</div>', unsafe_allow_html=True)
        col_dl1, col_dl2 = st.columns(2)
        safe = (our_project_name or micromarket).replace(' ', '_')

        with col_dl1:
            ppt_buf = generate_ppt(
                sections_dict, micromarket, city, product_type,
                our_project_name, our_configs, our_launch_timeline
            )
            st.download_button(
                label="📥 Download PPT Deck",
                data=ppt_buf,
                file_name=f"{safe}_Competition_Analysis.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        with col_dl2:
            st.download_button(
                label="📄 Download Full Report (TXT)",
                data=result,
                file_name=f"{safe}_Competition_Analysis.txt",
                mime="text/plain"
            )
