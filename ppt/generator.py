"""
PowerPoint slide generation for competition analysis reports.
"""
import io
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ppt.helpers import (
    add_bullets,
    add_footer,
    add_rectangle,
    add_textbox,
    clean_markdown,
    set_background,
    set_cell_background,
)
from ppt.theme import (
    CREAM,
    DARK_NAVY,
    GOLD,
    GOLD_LIGHT,
    GOLD_PALE,
    LIGHT_BLUE,
    MID_GRAY,
    NAVY,
    SLATE,
    WHITE,
    RED_FLAG,
    AMBER_FLAG,
    GREEN_FLAG,
)


def _cover(prs, micromarket, city, product_type, project_name, timeline):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(sl, NAVY)
    add_rectangle(sl, 0, 0, 0.12, 7.5, GOLD)
    add_rectangle(sl, 9.0, 0, 1.0, 0.14, GOLD)
    add_rectangle(sl, 9.0, 7.36, 1.0, 0.14, GOLD)
    add_textbox(sl, "REAL ESTATE DEVELOPER INTELLIGENCE", 0.5, 0.32, 9.0, 0.38, size=9, bold=True, color=GOLD)
    add_textbox(sl, "Competition\nAnalysis", 0.5, 1.05, 9.0, 1.9, size=52, bold=True, color=WHITE)
    add_textbox(sl, project_name or "Developer Project", 0.5, 3.1, 9.0, 0.6, size=24, color=GOLD_LIGHT)
    add_textbox(sl, f"{micromarket},  {city}", 0.5, 3.65, 9.0, 0.45, size=16, color=LIGHT_BLUE)
    add_rectangle(sl, 0.5, 4.3, 4.0, 0.05, GOLD)
    add_textbox(
        sl,
        f"Product Type: {product_type}    \u00b7    Launch: {timeline}",
        0.5, 4.48, 9.0, 0.38, size=11, color=LIGHT_BLUE,
    )
    add_textbox(
        sl,
        "STRICTLY CONFIDENTIAL  \u00b7  INTERNAL STRATEGY USE ONLY",
        0.5, 6.95, 9.0, 0.28, size=8, italic=True, color=MID_GRAY,
    )


def _exec_summary(prs, micromarket, city, bullets, project_name, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(sl, CREAM)
    add_rectangle(sl, 0, 0, 10, 0.88, NAVY)
    add_textbox(sl, "EXECUTIVE SUMMARY", 0.4, 0.2, 7.5, 0.5, size=20, bold=True, color=WHITE)
    add_textbox(sl, f"{micromarket}, {city}", 7.2, 0.23, 2.6, 0.4, size=10, color=GOLD_LIGHT, align=PP_ALIGN.RIGHT)
    labels = ["Market Snapshot", "Competitive Landscape", "Pricing Position", "Strategic Verdict"]
    icons = ["01", "02", "03", "04"]
    pos = [(0.28, 1.0), (5.18, 1.0), (0.28, 3.72), (5.18, 3.72)]
    for i, (lx, ly) in enumerate(pos):
        add_rectangle(sl, lx, ly, 4.6, 2.48, WHITE, line_color=RGBColor(0xD8, 0xD8, 0xE8), line_width=0.5)
        add_rectangle(sl, lx, ly, 4.6, 0.10, GOLD)
        add_textbox(sl, icons[i], lx + 0.16, ly + 0.14, 0.6, 0.42, size=24, bold=True, color=GOLD)
        add_textbox(sl, labels[i], lx + 0.74, ly + 0.17, 3.6, 0.35, size=10, bold=True, color=NAVY)
        raw = bullets[i] if i < len(bullets) else "See full analysis"
        add_textbox(sl, clean_markdown(raw), lx + 0.16, ly + 0.62, 4.2, 1.76, size=10, color=SLATE)
    add_footer(sl, f"{project_name or 'Project'}  \u00b7  {micromarket}, {city}  \u00b7  Confidential", page)


def _market(prs, micromarket, city, content, project_name, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(sl, CREAM)
    add_rectangle(sl, 0, 0, 10, 0.88, NAVY)
    add_textbox(sl, "MARKET SNAPSHOT", 0.4, 0.2, 7.5, 0.5, size=20, bold=True, color=WHITE)
    add_textbox(sl, f"{micromarket}, {city}", 7.2, 0.23, 2.6, 0.4, size=10, color=GOLD_LIGHT, align=PP_ALIGN.RIGHT)
    add_rectangle(sl, 0.28, 1.0, 3.0, 6.1, NAVY)
    pm = re.findall(r"\u20b9[\d,]+(?:\s*[\u2013-]\s*\u20b9[\d,]+)?", content)
    sv = [pm[0][:18] if pm else "\u20b918K\u201324K", "STABLE \u2192", "+8\u201312%", "~12 active"]
    sl_ = ["Avg \u20b9/sqft", "Market Temp", "YoY Growth", "Inventory"]
    for i, (lb, vl) in enumerate(zip(sl_, sv)):
        y = 1.3 + i * 1.35
        add_textbox(sl, lb, 0.42, y, 2.7, 0.28, size=8, bold=True, color=GOLD)
        add_textbox(sl, vl, 0.42, y + 0.29, 2.7, 0.62, size=14, bold=True, color=WHITE)
        if i < 3:
            add_rectangle(sl, 0.45, y + 0.98, 2.5, 0.02, DARK_NAVY)
    add_textbox(sl, "Market Intelligence", 3.55, 1.05, 6.1, 0.4, size=13, bold=True, color=NAVY)
    add_rectangle(sl, 3.55, 1.5, 1.6, 0.05, GOLD)
    lines = [l.strip().lstrip("*- ") for l in content.split("\n") if l.strip() and len(l.strip()) > 12]
    add_bullets(sl, lines, 3.55, 1.65, 6.1, row_height=0.55, max_rows=9)
    add_footer(sl, f"{project_name or 'Project'}  \u00b7  {micromarket}, {city}  \u00b7  Confidential", page)


def _deepdive(prs, micromarket, city, content, project_name, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(sl, CREAM)
    add_rectangle(sl, 0, 0, 10, 0.88, NAVY)
    add_textbox(sl, "COMPETITOR DEEP-DIVE", 0.4, 0.2, 7.5, 0.5, size=20, bold=True, color=WHITE)
    add_textbox(sl, f"{micromarket}, {city}", 7.2, 0.23, 2.6, 0.4, size=10, color=GOLD_LIGHT, align=PP_ALIGN.RIGHT)
    blocks = re.split(r"\*\*(.+?)\*\*", content)
    cards = []
    for j in range(1, len(blocks) - 1, 2):
        h = blocks[j].strip()
        b = blocks[j + 1].strip() if j + 1 < len(blocks) else ""
        if "\u2014" in h or "-" in h or len(h) > 5:
            cards.append((h, b))
    pos = [(0.28, 1.0), (5.18, 1.0), (0.28, 3.8), (5.18, 3.8)]
    if cards:
        for idx, (header, body) in enumerate(cards[:4]):
            lx, ly = pos[idx]
            add_rectangle(sl, lx, ly, 4.6, 2.55, WHITE, line_color=RGBColor(0xD8, 0xD8, 0xE8), line_width=0.5)
            add_rectangle(sl, lx, ly, 0.09, 2.55, GOLD)
            add_textbox(sl, header[:55], lx + 0.2, ly + 0.1, 4.2, 0.36, size=10, bold=True, color=NAVY)
            sub = [l.strip().lstrip("*- ") for l in body.split("\n") if l.strip() and len(l.strip()) > 5][:5]
            y = ly + 0.52
            for ln in sub:
                add_rectangle(sl, lx + 0.2, y + 0.1, 0.06, 0.06, GOLD)
                add_textbox(sl, ln[:120], lx + 0.38, y, 4.1, 0.4, size=9, color=SLATE)
                y += 0.4
    else:
        lines = [l.strip().lstrip("*- ") for l in content.split("\n") if l.strip() and len(l.strip()) > 8]
        add_bullets(sl, lines, 0.35, 1.05, 9.3, max_rows=14)
    add_footer(sl, f"{project_name or 'Project'}  \u00b7  {micromarket}, {city}  \u00b7  Confidential", page)


def _table(prs, micromarket, city, content, project_name, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(sl, CREAM)
    add_rectangle(sl, 0, 0, 10, 0.88, NAVY)
    add_textbox(sl, "COMPETITOR COMPARISON TABLE", 0.4, 0.2, 7.5, 0.5, size=20, bold=True, color=WHITE)
    add_textbox(sl, f"{micromarket}, {city}", 7.2, 0.23, 2.6, 0.4, size=10, color=GOLD_LIGHT, align=PP_ALIGN.RIGHT)
    rows = []
    for line in content.split("\n"):
        if "|" in line and "---" not in line and line.strip().startswith("|"):
            cells = [c.strip() for c in line.split("|") if c.strip()]
            if cells:
                rows.append(cells)
    if len(rows) >= 2:
        headers = rows[0]
        data_rows = rows[1:9]
        n_cols = min(len(headers), 8)
        n_rows = len(data_rows) + 1
        raw_w = [1.7, 1.4, 0.85, 1.0, 0.95, 1.1, 0.85, 1.7][:n_cols]
        scale = 9.3 / sum(raw_w)
        col_w = [w * scale for w in raw_w]
        rh = min(Inches(0.48), Inches(5.5 / n_rows))
        tbl = sl.shapes.add_table(n_rows, n_cols, Inches(0.35), Inches(1.05), Inches(9.3), rh * n_rows).table
        for ci, cw in enumerate(col_w):
            tbl.columns[ci].width = Inches(cw)
        for ci, h in enumerate(headers[:n_cols]):
            cell = tbl.cell(0, ci)
            cell.text = h.upper()
            set_cell_background(cell, 0x1E, 0x27, 0x61)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.font.bold = True
            run.font.size = Pt(8)
            run.font.color.rgb = GOLD
            run.font.name = "Calibri"
        for ri, row_data in enumerate(data_rows):
            bg = (0xFF, 0xFF, 0xFF) if ri % 2 == 0 else (0xF0, 0xF0, 0xF8)
            for ci in range(n_cols):
                cell = tbl.cell(ri + 1, ci)
                cell.text = row_data[ci] if ci < len(row_data) else ""
                set_cell_background(cell, *bg)
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.runs[0] if p.runs else p.add_run()
                run.font.size = Pt(9)
                run.font.bold = ci == 0
                run.font.color.rgb = NAVY if ci == 0 else SLATE
                run.font.name = "Calibri"
    else:
        add_textbox(
            sl,
            "Table data unavailable - see Competitor Deep-Dive tab.",
            0.35, 1.1, 9.3, 0.4, size=11, color=SLATE,
        )
    add_footer(sl, f"{project_name or 'Project'}  \u00b7  {micromarket}, {city}  \u00b7  Confidential", page)


def _pricing(prs, micromarket, city, content, project_name, configs_str, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(sl, CREAM)
    add_rectangle(sl, 0, 0, 10, 0.88, NAVY)
    add_textbox(
        sl,
        f"PRICING STRATEGY  \u00b7  {(project_name or 'YOUR PROJECT').upper()}",
        0.4, 0.2, 9.3, 0.5, size=18, bold=True, color=WHITE,
    )
    add_textbox(sl, f"{micromarket}, {city}", 7.2, 0.23, 2.6, 0.4, size=10, color=GOLD_LIGHT, align=PP_ALIGN.RIGHT)
    raw_c = [c.strip() for c in configs_str.split(",") if c.strip()] if configs_str else []
    defs = ["2BHK", "3BHK", "Premium / Jodi"]
    cfgs = [raw_c[i] if i < len(raw_c) else defs[i] for i in range(3)]
    pl = [
        l
        for l in content.split("\n")
        if ("\u2192" in l or "all-in" in l.lower() or "carpet" in l.lower()) and "\u20b9" in l
    ]
    prices = []
    for p in pl:
        m = re.search(r"\u20b9\d+\.\d+\s*Cr", p)
        if m:
            prices.append(m.group().strip())
    if not prices:
        prices = re.findall(r"\u20b9\d+\.\d+\s*Cr", content)
    cf = [NAVY, GOLD, NAVY]
    vc = [GOLD_LIGHT, WHITE, GOLD_LIGHT]
    sc = [LIGHT_BLUE, DARK_NAVY, LIGHT_BLUE]
    for i in range(3):
        lx = 0.28 + i * 3.25
        add_rectangle(sl, lx, 1.0, 3.1, 2.0, cf[i])
        add_textbox(sl, cfgs[i], lx + 0.2, 1.1, 2.7, 0.5, size=22, bold=True, color=WHITE)
        add_textbox(sl, prices[i] if i < len(prices) else "See report", lx + 0.2, 1.65, 2.7, 0.5, size=17, bold=True, color=vc[i])
        add_textbox(sl, "All-in price", lx + 0.2, 2.22, 2.7, 0.3, size=9, color=sc[i])
    add_textbox(sl, "Pricing Intelligence", 0.28, 3.2, 9.5, 0.38, size=13, bold=True, color=NAVY)
    add_rectangle(sl, 0.28, 3.62, 1.6, 0.05, GOLD)
    lines = [l.strip().lstrip("*- ") for l in content.split("\n") if l.strip() and len(l.strip()) > 8]
    add_bullets(sl, lines, 0.28, 3.78, 9.5, row_height=0.46, max_rows=8)
    add_footer(sl, f"{project_name or 'Project'}  \u00b7  {micromarket}, {city}  \u00b7  Confidential", page)


def _gaps(prs, micromarket, city, content, project_name, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(sl, CREAM)
    add_rectangle(sl, 0, 0, 10, 0.88, NAVY)
    add_textbox(sl, "MARKET GAPS & DIFFERENTIATION OPPORTUNITIES", 0.4, 0.2, 9.3, 0.5, size=18, bold=True, color=WHITE)
    add_textbox(sl, f"{micromarket}, {city}", 7.2, 0.23, 2.6, 0.4, size=10, color=GOLD_LIGHT, align=PP_ALIGN.RIGHT)
    lines = [l.strip().lstrip("*- ") for l in content.split("\n") if l.strip() and len(l.strip()) > 10]
    gap_lines = [l for l in lines if re.match(r"(?i)gap", l)][:3]
    if not gap_lines:
        gap_lines = lines[:3]
    add_textbox(sl, "IDENTIFIED GAPS", 0.28, 0.98, 4.5, 0.3, size=9, bold=True, color=NAVY)
    for i, gap in enumerate(gap_lines):
        c2 = re.sub(r"(?i)^gap\s*\d*\s*[:.]?\s*", "", gap).strip()
        ly = 1.35 + i * 1.75
        add_rectangle(sl, 0.28, ly, 4.5, 1.56, WHITE, line_color=RGBColor(0xD8, 0xD8, 0xE8), line_width=0.5)
        add_rectangle(sl, 0.28, ly, 0.09, 1.56, GOLD)
        add_textbox(sl, f"GAP {i + 1:02d}", 0.46, ly + 0.1, 4.0, 0.28, size=9, bold=True, color=GOLD)
        add_textbox(sl, c2[:220], 0.46, ly + 0.42, 3.95, 1.02, size=10, color=SLATE)
    add_rectangle(sl, 5.02, 0.98, 4.72, 6.1, NAVY)
    add_textbox(sl, "HOW TO WIN", 5.22, 1.12, 4.2, 0.32, size=9, bold=True, color=GOLD)
    add_rectangle(sl, 5.22, 1.48, 1.2, 0.04, GOLD_PALE)
    win = []
    for l in lines:
        cl = re.sub(r"(?i)(how\s+\S.{0,30}can\s+win\s*[:\u00b7]?)", "", l).strip().lstrip(":\u00b7- ")
        if cl and len(cl) >= 8 and l not in gap_lines:
            win.append(cl)
    add_bullets(sl, win, 5.22, 1.65, 4.35, row_height=0.56, max_rows=8, dot_color=GOLD, text_color=LIGHT_BLUE)
    add_footer(sl, f"{project_name or 'Project'}  \u00b7  {micromarket}, {city}  \u00b7  Confidential", page)


def _risks(prs, micromarket, city, content, project_name, page):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(sl, CREAM)
    add_rectangle(sl, 0, 0, 10, 0.88, NAVY)
    add_textbox(sl, "RISK FLAGS", 0.4, 0.2, 7.5, 0.5, size=20, bold=True, color=WHITE)
    add_textbox(sl, f"{micromarket}, {city}", 7.2, 0.23, 2.6, 0.4, size=10, color=GOLD_LIGHT, align=PP_ALIGN.RIGHT)
    lines = [l.strip().lstrip("*- ") for l in content.split("\n") if l.strip() and len(l.strip()) > 8]
    verdict = ""
    y = 1.02
    count = 0
    cmap = {"HIGH": RED_FLAG, "MED": AMBER_FLAG, "LOW": GREEN_FLAG}
    for line in lines:
        if "VERDICT" in line.upper():
            verdict = line
            continue
        lvl = "MED"
        for lv in ["HIGH", "MED", "LOW"]:
            if f"[{lv}]" in line.upper() or f" {lv}" in line.upper():
                lvl = lv
                break
        add_rectangle(sl, 0.28, y, 1.05, 0.72, cmap[lvl])
        add_textbox(sl, lvl, 0.28, y + 0.17, 1.05, 0.4, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rectangle(sl, 1.45, y, 8.25, 0.72, WHITE, line_color=RGBColor(0xD8, 0xD8, 0xE8), line_width=0.5)
        c2 = re.sub(r"\[?(HIGH|MED|LOW)\]?", "", line, flags=re.IGNORECASE).strip().lstrip(":- ")
        add_textbox(sl, c2[:220], 1.6, y + 0.08, 7.95, 0.58, size=10, color=SLATE)
        y += 0.88
        count += 1
        if count >= 5:
            break
    if verdict:
        add_rectangle(sl, 0.28, y + 0.18, 9.44, 0.88, NAVY)
        vc = re.sub(r"(?i)^verdict\s*[:\u00b7]?\s*", "", verdict).strip()
        add_textbox(sl, f"VERDICT  \u00b7  {vc}", 0.46, y + 0.3, 9.1, 0.64, size=12, bold=True, color=GOLD_LIGHT)
    add_footer(sl, f"{project_name or 'Project'}  \u00b7  {micromarket}, {city}  \u00b7  Confidential", page)


def _closing(prs, micromarket, city, project_name):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(sl, NAVY)
    add_rectangle(sl, 0, 0, 0.12, 7.5, GOLD)
    add_rectangle(sl, 9.0, 0, 1.0, 0.14, GOLD)
    add_rectangle(sl, 9.0, 7.36, 1.0, 0.14, GOLD)
    add_textbox(sl, "REAL ESTATE DEVELOPER INTELLIGENCE", 0.5, 1.9, 9.0, 0.45, size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(sl, "Thank You", 0.5, 2.7, 9.0, 1.1, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rectangle(sl, 3.0, 4.0, 4.0, 0.07, GOLD)
    add_textbox(
        sl,
        f"{project_name or 'Developer Project'}  \u00b7  {micromarket}, {city}",
        0.5, 4.2, 9.0, 0.48, size=15, color=LIGHT_BLUE, align=PP_ALIGN.CENTER,
    )
    add_textbox(
        sl,
        "This document is confidential and intended solely for internal strategy use.",
        0.5, 6.7, 9.0, 0.35, size=9, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER,
    )


def generate_ppt(
    sections: dict[str, str],
    micromarket: str,
    city: str,
    product_type: str,
    project_name: str,
    configs_str: str,
    launch_timeline: str,
) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    eb = [sections.get(k, "")[:200] for k in ["1", "2", "4", "6"]]
    _cover(prs, micromarket, city, product_type, project_name, launch_timeline)
    _exec_summary(prs, micromarket, city, eb, project_name, 2)
    _market(prs, micromarket, city, sections.get("1", ""), project_name, 3)
    _deepdive(prs, micromarket, city, sections.get("2", ""), project_name, 4)
    _table(prs, micromarket, city, sections.get("3", ""), project_name, 5)
    _pricing(prs, micromarket, city, sections.get("4", ""), project_name, configs_str, 6)
    _gaps(prs, micromarket, city, sections.get("5", ""), project_name, 7)
    _risks(prs, micromarket, city, sections.get("6", ""), project_name, 8)
    _closing(prs, micromarket, city, project_name)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
