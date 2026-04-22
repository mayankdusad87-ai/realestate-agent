"""
Low-level PowerPoint shape and text helpers.
"""
import re

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from ppt.theme import GOLD, LIGHT_BLUE, NAVY, SLATE, WHITE


def set_background(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_rectangle(slide, left, top, width, height, fill_color, line_color=None, line_width=0.5):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    text,
    left,
    top,
    width,
    height,
    size=11,
    bold=False,
    italic=False,
    color=None,
    align=PP_ALIGN.LEFT,
    font="Calibri",
):
    if color is None:
        color = WHITE
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def set_cell_background(cell, r: int, g: int, b: int):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    sf = etree.SubElement(tcPr, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", f"{r:02X}{g:02X}{b:02X}")


def add_footer(slide, label: str, page: int):
    add_rectangle(slide, 0, 7.28, 10, 0.22, NAVY)
    add_textbox(slide, label, 0.3, 7.30, 8.8, 0.19, size=8, color=LIGHT_BLUE)
    add_textbox(
        slide, str(page), 9.5, 7.30, 0.4, 0.19, size=8, color=GOLD, align=PP_ALIGN.RIGHT
    )


def add_bullets(slide, lines, start_x, start_y, max_width, row_height=0.44, max_rows=10, dot_color=None, text_color=None):
    if dot_color is None:
        dot_color = GOLD
    if text_color is None:
        text_color = SLATE
    y = start_y
    for line in lines[:max_rows]:
        cleaned = line.strip().lstrip("*-\t ")
        if not cleaned or len(cleaned) < 4:
            continue
        add_rectangle(slide, start_x, y + 0.15, 0.06, 0.06, dot_color)
        add_textbox(slide, cleaned, start_x + 0.18, y, max_width - 0.18, row_height, size=10, color=text_color)
        y += row_height
        if y > 7.0:
            break
    return y


def clean_markdown(text: str, max_len: int = 260) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    lines = [
        l.strip().lstrip("-*\t ")
        for l in text.split("\n")
        if l.strip() and len(l.strip()) > 5
    ]
    return ("  ·  ".join(lines[:3]))[:max_len]
